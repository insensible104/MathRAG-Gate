# gen_ppt.py
"""
Script to generate an editable PowerPoint (.pptx) presentation for MathRAG-Gate.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def find_latest_benchmark_image(results_dir="results"):
    """Finds the benchmark_results.png from the latest experiment folder."""
    if not os.path.exists(results_dir):
        return None
    
    # Get all subdirectories starting with 'exp_'
    exp_dirs = [os.path.join(results_dir, d) for d in os.listdir(results_dir) 
                if os.path.isdir(os.path.join(results_dir, d)) and d.startswith("exp_")]
    
    if not exp_dirs:
        return None
        
    # Sort by creation time (latest first)
    latest_dir = max(exp_dirs, key=os.path.getmtime)
    image_path = os.path.join(latest_dir, "benchmark_results.png")
    
    if os.path.exists(image_path):
        print(f"✅ Found benchmark image at: {image_path}")
        return image_path
    return None

def create_presentation():
    prs = Presentation()

    # Helper to set title and content
    def add_slide(layout_idx, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if content_text and layout_idx == 1: # Title and Content
            content = slide.placeholders[1]
            content.text = content_text
        return slide

    # --- Slide 1: Title ---
    slide = add_slide(0, "MathRAG-Gate") # Title Slide
    subtitle = slide.placeholders[1]
    subtitle.text = "A Confidence-Gated RAG Framework for\nRobust Mathematical Reasoning\n\nDSAA 5020 Course Project"

    # --- Slide 2: The Problem ---
    content = (
        "1. Relevance ≠ Quality\n"
        "   - Traditional RAG optimizes for semantic relevance, often retrieving documents with correct numbers but wrong logic.\n\n"
        "2. Fusion Noise\n"
        "   - Standard Hybrid Retrieval (RRF) introduces low-quality documents, causing performance degradation.\n\n"
        "3. Engineering Bottleneck\n"
        "   - Running dual models (7B Gen + 0.5B Judge) on 12GB VRAM causes severe thrashing."
    )
    add_slide(1, "The Problem: The Logical Quality Gap", content)

    # --- Slide 3: Solution Architecture ---
    content = (
        "1. Confidence Gate (Core Innovation)\n"
        "   - Dynamically chooses between Rule-RQP (Fast) and LLM-RQP (Robust) based on consistency check (rho).\n\n"
        "2. Staged Batching (Engineering)\n"
        "   - Decouples Retrieval and Generation phases.\n"
        "   - Reduces model swapping from 500 times to 1 time.\n\n"
        "3. CoT Provider\n"
        "   - Acts as a logic supplier for the Generator's Chain-of-Thought."
    )
    add_slide(1, "Solution: MathRAG-Gate Architecture", content)

    # --- Slide 4: Methodology - The Judge ---
    content = (
        "Rule-RQP (Efficiency)\n"
        " - Detects LaTeX structure, causal markers ('Therefore'), and boxed answers.\n\n"
        "LLM-RQP (Robustness)\n"
        " - Powered by Qwen2.5-0.5B with Few-Shot Prompting to stabilize scoring.\n"
        " - Used when Rule-based signals are ambiguous."
    )
    add_slide(1, "Methodology: Reasoning Quality Predictors", content)

    # --- Slide 5: Engineering Optimization ---
    content = (
        "Challenge: 12GB VRAM Constraint\n"
        " - Concurrent execution caused >13s latency and OOM errors.\n\n"
        "Solution: Staged Batching\n"
        " - Phase 1: Batch Retrieval & Scoring (Only 0.5B loaded).\n"
        " - Phase 2: Batch Generation (Only 7B loaded).\n\n"
        "Result: Latency stabilized at ~4.5s with 0 errors."
    )
    add_slide(1, "Engineering: Staged Batching Strategy", content)

    # --- Slide 6: Experimental Setup ---
    content = (
        "Dataset: MATH Dataset (Competition Math)\n"
        "Hardware: NVIDIA RTX 4070 Ti (12GB)\n"
        "Models: Qwen2.5-7B (Generator), Qwen2.5-0.5B (Judge)\n"
        "Baselines: Dense (BGE), Sparse (BM25), Hybrid (RRF)"
    )
    add_slide(1, "Experimental Setup", content)

    # --- Slide 7: Results - Fusion Noise ---
    slide = add_slide(1, "Results: Discovery of Fusion Noise", "")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Baseline Performance:"
    p = tf.add_paragraph()
    p.text = "• Dense Retrieval: 86.00%"
    p = tf.add_paragraph()
    p.text = "• Sparse Retrieval: 85.00%"
    p = tf.add_paragraph()
    p.text = "• Hybrid (RRF): 81.00% (Performance Drop!)"
    p = tf.add_paragraph()
    p.text = "\nInsight: Naive fusion introduces noise, degrading performance below the single-model baseline."

    # Try to insert image if available
    img_path = find_latest_benchmark_image()
    if img_path:
        # Add slide for visual
        slide_img = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
        title = slide_img.shapes.title
        title.text = "Benchmark Visualization"
        slide_img.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

    # --- Slide 8: Results - Performance Recovery ---
    content = (
        "MathRAG-Gate Performance: 86.80%\n\n"
        "1. Recovery: Successfully recovered from Hybrid's low point (81.00%) to surpass Dense (86.00%).\n\n"
        "2. Stability: Accuracy Std Dev is only 0.0110 across 5 runs.\n\n"
        "3. Conclusion: The Gate acts as a 'Safety Layer', filtering out low-quality documents."
    )
    add_slide(1, "Results: MathRAG-Gate Recovery", content)

    # --- Slide 9: Hyperparameter Insight ---
    content = (
        "Grid Search Result:\n"
        "Optimal Configuration: RRF Weight = 0.1 / Quality Weight = 0.9\n\n"
        "Interpretation:\n"
        "The system performs best when it almost entirely trusts the internal Judge's quality assessment over the external retrieval rank.\n"
        "This validates the effectiveness of our Few-Shot Optimized 0.5B Judge."
    )
    add_slide(1, "Insight: Trusting the Judge", content)

    # --- Slide 10: The Role of CoT ---
    content = (
        "Initial Accuracy (No CoT): ~30%\n"
        "Final Accuracy (With CoT): ~87%\n\n"
        "Synergy:\n"
        "Recall was always high (99%), but the Generator couldn't use it.\n"
        "MathRAG-Gate acts as a 'CoT Enabler' by providing logically sound templates that the generator can imitate."
    )
    add_slide(1, "The Impact of Chain-of-Thought", content)

    # --- Slide 11: Conclusion ---
    content = (
        "1. Academic Contribution\n"
        "   - Transformed RAG from 'Semantic Search' to 'Logical Quality Assurance'.\n"
        "   - Proved that Quality > Relevance for math tasks.\n\n"
        "2. Engineering Achievement\n"
        "   - Solved VRAM bottleneck on consumer GPU.\n"
        "   - Achieved stable high-throughput inference."
    )
    add_slide(1, "Conclusion", content)

    # Save
    output_file = "MathRAG-Gate_Presentation.pptx"
    prs.save(output_file)
    print(f"🎉 Presentation saved as: {os.path.abspath(output_file)}")

if __name__ == "__main__":
    create_presentation()